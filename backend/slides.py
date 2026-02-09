
import os
from pptx import Presentation
from pptx.util import Pt

def create_presentation(slides_data, output_path):
    try:
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0] 
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "ScholarMate Analysis"
        subtitle.text = "AI-Generated Research Report"
        
        content_slide_layout = prs.slide_layouts[1] 
        
        for slide_info in slides_data:
            slide = prs.slides.add_slide(content_slide_layout)
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = slide_info.get('title', 'Untitled Slide')
            
            body_shape = slide.placeholders[1]
            if body_shape:
                tf = body_shape.text_frame
                tf.word_wrap = True
                bullets = slide_info.get('bullets', [])
                if bullets:
                    tf.text = bullets[0]
                    for bullet_text in bullets[1:]:
                        p = tf.add_paragraph()
                        p.text = bullet_text
                        p.level = 0
                        p.space_before = Pt(6)

        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        prs.save(output_path)
        return True
    except Exception as e:
        print(f"Error creating presentation: {e}")
        return False
